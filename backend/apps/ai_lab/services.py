import ast
import csv
import html as html_lib
import json
import logging
import os
import random
import re
import zipfile
from io import BytesIO
from pathlib import Path
from urllib import error as urllib_error
from urllib import request as urllib_request
from xml.etree import ElementTree

from apps.administration.utils import setting_value

logger = logging.getLogger(__name__)

GROQ_MODEL = os.getenv("GROQ_MODEL", "llama-3.3-70b-versatile")
OPENROUTER_MODELS = [
    model.strip()
    for model in os.getenv(
        "OPENROUTER_MODELS",
        "nvidia/nemotron-3-nano-30b-a3b:free,google/gemma-4-31b-it:free,openrouter/free",
    ).split(",")
    if model.strip()
]
SUPPORTED_TEXT_EXTENSIONS = {
    ".pdf",
    ".docx",
    ".pptx",
    ".xlsx",
    ".xlsm",
    ".csv",
    ".txt",
    ".md",
    ".json",
    ".xml",
    ".html",
    ".htm",
    ".rtf",
    ".odt",
}
AI_CONTEXT_LIMIT = 25000
AI_TOPIC_CONTEXT_LIMIT = 10000
AI_RETRY_ATTEMPTS = 2
TOPIC_STOPWORDS = {
    "about",
    "after",
    "also",
    "because",
    "before",
    "between",
    "chapter",
    "course",
    "data",
    "document",
    "example",
    "from",
    "have",
    "lecture",
    "model",
    "more",
    "other",
    "section",
    "should",
    "student",
    "system",
    "than",
    "that",
    "their",
    "there",
    "these",
    "this",
    "using",
    "value",
    "when",
    "which",
    "with",
    "will",
}

def extract_document_text(file_path):
    path = Path(file_path)
    try:
        return extract_document_text_from_bytes(path.read_bytes(), path.name)
    except Exception as exc:
        logger.warning("Unable to read document %s: %s", file_path, exc)
        return ""

def extract_document_text_from_upload(uploaded_file):
    name = getattr(uploaded_file, "name", "") or "uploaded-file"
    if hasattr(uploaded_file, "chunks"):
        data = b"".join(chunk for chunk in uploaded_file.chunks())
    else:
        data = uploaded_file.read()
    return extract_document_text_from_bytes(data, name)


def extract_document_text_from_bytes(data, file_name):
    extension = Path(file_name).suffix.lower()

    if extension == ".pdf":
        return clean_extracted_text(extract_pdf_text(BytesIO(data)))
    if extension == ".docx":
        return clean_extracted_text(extract_docx_text(BytesIO(data)))
    if extension == ".pptx":
        return clean_extracted_text(extract_pptx_text(BytesIO(data)))
    if extension in {".xlsx", ".xlsm"}:
        return clean_extracted_text(extract_xlsx_text(BytesIO(data)))
    if extension == ".csv":
        return clean_extracted_text(extract_csv_text(data))
    if extension in {".html", ".htm"}:
        return clean_extracted_text(extract_html_text(data))
    if extension == ".json":
        return clean_extracted_text(extract_json_text(data))
    if extension == ".xml":
        return clean_extracted_text(extract_xml_text(data))
    if extension == ".rtf":
        return clean_extracted_text(extract_rtf_text(data))
    if extension == ".odt":
        return clean_extracted_text(extract_odt_text(BytesIO(data)))
    if extension in {".txt", ".md"}:
        return clean_extracted_text(extract_plain_text_from_bytes(data))

    return ""


def read_binary(file_obj):
    if isinstance(file_obj, (bytes, bytearray)):
        return bytes(file_obj)
    if hasattr(file_obj, "getvalue"):
        return file_obj.getvalue()
    if hasattr(file_obj, "seek"):
        file_obj.seek(0)
    data = file_obj.read() if hasattr(file_obj, "read") else Path(file_obj).read_bytes()
    if hasattr(file_obj, "seek"):
        file_obj.seek(0)
    return data


def extract_pdf_text(file_obj):
    data = read_binary(file_obj)
    rich_text = extract_pdf_text_with_pdfplumber(data)
    if rich_text:
        return rich_text

    text_parts = []
    try:
        from pypdf import PdfReader

        reader = PdfReader(BytesIO(data))
        for index, page in enumerate(reader.pages, start=1):
            page_text = page.extract_text() or ""
            if page_text.strip():
                text_parts.append(f"Page {index}\nText\n{page_text.strip()}")
    except Exception as exc:
        logger.warning("Unable to extract PDF text: %s", exc)

    fallback_text = "\n\n".join(text_parts).strip()
    ocr_text = extract_pdf_ocr_text(data)
    return "\n\n".join(part for part in (fallback_text, ocr_text) if part).strip()


def extract_pdf_text_with_pdfplumber(data):
    try:
        import pdfplumber
    except Exception as exc:
        logger.debug("pdfplumber extraction unavailable: %s", exc)
        return ""

    pages = []
    try:
        with pdfplumber.open(BytesIO(data)) as pdf:
            metadata = []
            if pdf.metadata:
                for key in ("Title", "Author", "Subject", "Creator", "Producer"):
                    value = pdf.metadata.get(key)
                    if value:
                        metadata.append(f"{key}: {value}")
            if metadata:
                pages.append("Document metadata\n" + "\n".join(metadata))

            for page_index, page in enumerate(pdf.pages, start=1):
                page_parts = [f"Page {page_index}"]
                text = extract_pdf_layout_text(page)
                if text:
                    page_parts.append("Layout text\n" + text)

                tables = extract_pdf_tables(page)
                if tables:
                    page_parts.append("Tables\n" + tables)

                needs_ocr = len(text or "") < 80
                ocr_text = extract_pdf_page_ocr_text(data, page_index - 1) if needs_ocr else ""
                if ocr_text:
                    page_parts.append("OCR text\n" + ocr_text)

                if len(page_parts) > 1:
                    pages.append("\n".join(page_parts))
    except Exception as exc:
        logger.warning("pdfplumber extraction failed: %s", exc)
        return ""

    return "\n\n".join(pages).strip()


def extract_pdf_layout_text(page):
    try:
        words = page.extract_words(
            keep_blank_chars=False,
            use_text_flow=True,
            extra_attrs=["size", "fontname"],
        )
    except Exception:
        raw_text = page.extract_text() or ""
        return mark_probable_headings(raw_text)

    if not words:
        return ""

    sizes = [float(word.get("size") or 0) for word in words if word.get("size")]
    body_size = sorted(sizes)[len(sizes) // 2] if sizes else 0
    lines = []
    current_top = None
    current_words = []

    for word in words:
        top = round(float(word.get("top", 0)), 1)
        if current_top is None or abs(top - current_top) <= 3:
            current_words.append(word)
            current_top = top if current_top is None else current_top
            continue
        lines.append(format_pdf_line(current_words, body_size))
        current_words = [word]
        current_top = top

    if current_words:
        lines.append(format_pdf_line(current_words, body_size))

    return "\n".join(line for line in lines if line).strip()


def format_pdf_line(words, body_size):
    sorted_words = sorted(words, key=lambda word: float(word.get("x0", 0)))
    pieces = []
    previous_x1 = None
    for word in sorted_words:
        word_text = word.get("text", "").strip()
        if not word_text:
            continue
        x0 = float(word.get("x0", 0))
        if previous_x1 is not None and x0 - previous_x1 > 28:
            pieces.append("|")
        pieces.append(word_text)
        previous_x1 = float(word.get("x1", x0))
    text = " ".join(pieces).replace(" | ", " | ").strip()
    if not text:
        return ""
    sizes = [float(word.get("size") or 0) for word in sorted_words if word.get("size")]
    avg_size = sum(sizes) / len(sizes) if sizes else 0
    is_heading = (
        avg_size and body_size and avg_size >= body_size * 1.18
    ) or (
        len(text) <= 90 and text == text.upper() and any(char.isalpha() for char in text)
    )
    return f"Heading: {text}" if is_heading else text


def mark_probable_headings(text):
    lines = []
    for line in (text or "").splitlines():
        stripped = line.strip()
        if not stripped:
            continue
        if len(stripped) <= 90 and (stripped == stripped.upper() or re.match(r"^\d+(\.\d+)*\s+\S+", stripped)):
            lines.append(f"Heading: {stripped}")
        else:
            lines.append(stripped)
    return "\n".join(lines)


def extract_pdf_tables(page):
    table_blocks = []
    try:
        tables = page.extract_tables() or []
    except Exception as exc:
        logger.debug("PDF table extraction failed on page: %s", exc)
        return ""

    for table_index, table in enumerate(tables, start=1):
        rows = []
        for row in table:
            values = [clean_cell(value) for value in row if clean_cell(value)]
            if values:
                rows.append(" | ".join(values))
        if rows:
            table_blocks.append(f"Table {table_index}\n" + "\n".join(rows))
    return "\n\n".join(table_blocks)


def clean_cell(value):
    return re.sub(r"\s+", " ", str(value or "")).strip()


def ocr_image(image):
    try:
        import pytesseract

        text = pytesseract.image_to_string(image) or ""
        return text.strip()
    except Exception as exc:
        logger.debug("OCR unavailable or failed: %s", exc)
        return ""


def extract_pdf_page_ocr_text(data, page_index):
    try:
        import fitz

        document = fitz.open(stream=data, filetype="pdf")
        page = document.load_page(page_index)
        pixmap = page.get_pixmap(matrix=fitz.Matrix(2, 2), alpha=False)
        from PIL import Image

        image = Image.open(BytesIO(pixmap.tobytes("png")))
        return ocr_image(image)
    except Exception as exc:
        logger.debug("PDF page OCR failed: %s", exc)
        return ""


def extract_pdf_ocr_text(data, max_pages=6):
    blocks = []
    try:
        import fitz

        document = fitz.open(stream=data, filetype="pdf")
        page_count = min(len(document), max_pages)
    except Exception as exc:
        logger.debug("Unable to open PDF for OCR: %s", exc)
        return ""

    for page_index in range(page_count):
        text = extract_pdf_page_ocr_text(data, page_index)
        if text:
            blocks.append(f"Page {page_index + 1} OCR text\n{text}")
    return "\n\n".join(blocks)


def clean_extracted_text(text):
    text = str(text or "").replace("\x00", " ")
    lines = []
    for raw_line in text.replace("\r\n", "\n").replace("\r", "\n").split("\n"):
        line = re.sub(r"[ \t\f\v]+", " ", raw_line).strip()
        if line:
            lines.append(line)
        elif lines and lines[-1] != "":
            lines.append("")
    cleaned = "\n".join(lines)
    cleaned = re.sub(r"\n{3,}", "\n\n", cleaned)
    return cleaned.strip()


def extract_plain_text_from_bytes(data):
    for encoding in ("utf-8-sig", "utf-8", "utf-16", "utf-16-le", "utf-16-be", "cp1252", "latin-1"):
        try:
            return data.decode(encoding).strip()
        except UnicodeDecodeError:
            continue
    return ""


def extract_plain_text(file_path):
    for encoding in ("utf-8", "utf-16", "latin-1"):
        try:
            return Path(file_path).read_text(encoding=encoding).strip()
        except UnicodeDecodeError:
            continue
        except Exception as exc:
            logger.warning("Unable to extract plain text from %s: %s", file_path, exc)
            return ""

    return ""


def extract_docx_text(file_path):
    library_text = extract_docx_text_with_library(file_path)
    if library_text:
        return library_text

    try:
        with zipfile.ZipFile(file_path) as archive:
            parts = []
            for name in (
                "word/document.xml",
                "word/footnotes.xml",
                "word/endnotes.xml",
                "word/comments.xml",
            ):
                try:
                    text = extract_openxml_text(archive.read(name))
                    if text:
                        parts.append(text)
                except KeyError:
                    continue
    except Exception as exc:
        logger.warning("Unable to read DOCX file %s: %s", file_path, exc)
        return ""

    return "\n\n".join(parts)


def extract_pptx_text(file_path):
    library_text = extract_pptx_text_with_library(file_path)
    if library_text:
        return library_text

    text_parts = []

    try:
        with zipfile.ZipFile(file_path) as archive:
            slide_names = sorted(name for name in archive.namelist() if name.startswith("ppt/slides/slide") and name.endswith(".xml"))
            for slide_name in slide_names:
                text = extract_openxml_text(archive.read(slide_name))
                if text:
                    text_parts.append(text)
    except Exception as exc:
        logger.warning("Unable to read PPTX file %s: %s", file_path, exc)
        return ""

    return "\n\n".join(text_parts).strip()


def extract_docx_text_with_library(file_obj):
    try:
        from docx import Document

        document = Document(file_obj)
        parts = []
        properties = document.core_properties
        metadata = []
        for label, value in (
            ("Title", properties.title),
            ("Author", properties.author),
            ("Subject", properties.subject),
            ("Keywords", properties.keywords),
            ("Category", properties.category),
        ):
            if value:
                metadata.append(f"{label}: {value}")
        if metadata:
            parts.append("Document metadata\n" + "\n".join(metadata))

        section = []
        for paragraph in document.paragraphs:
            text = paragraph.text.strip()
            if not text:
                continue
            style_name = getattr(paragraph.style, "name", "") or ""
            if style_name.lower().startswith("heading") or style_name.lower() in {"title", "subtitle"}:
                if section:
                    parts.append("\n".join(section))
                    section = []
                parts.append(f"{style_name}: {text}")
            else:
                section.append(text)
        if section:
            parts.append("\n".join(section))

        for table_index, table in enumerate(document.tables, start=1):
            rows = []
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                if cells:
                    rows.append(" | ".join(cells))
            if rows:
                parts.append(f"Table {table_index}\n" + "\n".join(rows))
        return "\n\n".join(parts).strip()
    except Exception as exc:
        logger.debug("python-docx extraction unavailable or failed: %s", exc)
        return ""


def extract_pptx_text_with_library(file_obj):
    try:
        from pptx import Presentation

        data = read_binary(file_obj)
        presentation = Presentation(BytesIO(data))
        notes_by_slide = extract_pptx_notes(data)
        slides = []
        for slide_index, slide in enumerate(presentation.slides, start=1):
            parts = [f"Slide {slide_index}"]
            title_shape = getattr(slide.shapes, "title", None)
            if title_shape and getattr(title_shape, "text", "").strip():
                parts.append(f"Title: {title_shape.text.strip()}")

            for shape in slide.shapes:
                if title_shape is not None and shape == title_shape:
                    continue
                if getattr(shape, "has_text_frame", False):
                    text = shape.text.strip()
                    if text:
                        parts.append("Text\n" + text)
                if getattr(shape, "has_table", False):
                    rows = []
                    for row in shape.table.rows:
                        cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                        if cells:
                            rows.append(" | ".join(cells))
                    if rows:
                        parts.append("Table\n" + "\n".join(rows))

                image_text = extract_pptx_shape_ocr(shape)
                if image_text:
                    parts.append("OCR text from slide image\n" + image_text)

            notes = notes_by_slide.get(slide_index, "")
            if notes:
                parts.append("Speaker notes\n" + notes)

            if len(parts) > 1:
                slides.append("\n".join(parts))
        return "\n\n".join(slides).strip()
    except Exception as exc:
        logger.debug("python-pptx extraction unavailable or failed: %s", exc)
        return ""


def extract_pptx_shape_ocr(shape):
    try:
        image = shape.image
    except Exception:
        return ""
    try:
        from PIL import Image

        pil_image = Image.open(BytesIO(image.blob))
        return ocr_image(pil_image)
    except Exception as exc:
        logger.debug("PPTX image OCR failed: %s", exc)
        return ""


def extract_pptx_notes(data):
    notes = {}
    try:
        with zipfile.ZipFile(BytesIO(data)) as archive:
            note_names = sorted(
                name
                for name in archive.namelist()
                if name.startswith("ppt/notesSlides/notesSlide") and name.endswith(".xml")
            )
            for index, name in enumerate(note_names, start=1):
                text = extract_openxml_text(archive.read(name))
                if text:
                    notes[index] = text
    except Exception as exc:
        logger.debug("PPTX notes extraction failed: %s", exc)
    return notes


def extract_xlsx_text(file_obj):
    try:
        from openpyxl import load_workbook

        workbook = load_workbook(file_obj, read_only=True, data_only=True)
        sheets = []
        for sheet in workbook.worksheets:
            rows = []
            for row in sheet.iter_rows(values_only=True):
                values = [str(value).strip() for value in row if value is not None and str(value).strip()]
                if values:
                    rows.append(" | ".join(values))
            if rows:
                sheets.append(f"Sheet: {sheet.title}\n" + "\n".join(rows))
        return "\n\n".join(sheets).strip()
    except Exception as exc:
        logger.warning("Unable to extract spreadsheet text: %s", exc)
        return ""


def extract_csv_text(data):
    text = extract_plain_text_from_bytes(data)
    if not text:
        return ""
    try:
        sample = text[:4096]
        dialect = csv.Sniffer().sniff(sample)
    except Exception:
        dialect = csv.excel
    rows = []
    try:
        for row in csv.reader(text.splitlines(), dialect):
            values = [value.strip() for value in row if value.strip()]
            if values:
                rows.append(" | ".join(values))
    except Exception:
        return text
    return "\n".join(rows)


def extract_html_text(data):
    text = extract_plain_text_from_bytes(data)
    if not text:
        return ""
    try:
        from bs4 import BeautifulSoup

        soup = BeautifulSoup(text, "html.parser")
        for element in soup(["script", "style", "noscript"]):
            element.decompose()
        return soup.get_text("\n")
    except Exception:
        without_scripts = re.sub(r"(?is)<(script|style).*?>.*?</\1>", " ", text)
        without_tags = re.sub(r"(?s)<[^>]+>", "\n", without_scripts)
        return html_lib.unescape(without_tags)


def extract_json_text(data):
    text = extract_plain_text_from_bytes(data)
    if not text:
        return ""
    try:
        parsed = json.loads(text)
    except Exception:
        return text
    return json.dumps(parsed, indent=2, ensure_ascii=False)


def extract_xml_text(data):
    text = extract_plain_text_from_bytes(data)
    if not text:
        return ""
    try:
        root = ElementTree.fromstring(text)
    except Exception:
        return re.sub(r"(?s)<[^>]+>", "\n", text)

    parts = []

    def walk(node, path=""):
        tag = node.tag.rsplit("}", 1)[-1]
        next_path = f"{path}/{tag}" if path else tag
        node_text = (node.text or "").strip()
        if node_text:
            parts.append(f"{next_path}: {node_text}")
        for child in node:
            walk(child, next_path)

    walk(root)
    return "\n".join(parts)


def extract_rtf_text(data):
    text = extract_plain_text_from_bytes(data)
    if not text:
        return ""
    text = re.sub(r"\\'[0-9a-fA-F]{2}", " ", text)
    text = re.sub(r"\\[a-zA-Z]+-?\d* ?", " ", text)
    text = text.replace("{", " ").replace("}", " ")
    return text


def extract_odt_text(file_obj):
    try:
        with zipfile.ZipFile(file_obj) as archive:
            xml = archive.read("content.xml")
    except Exception as exc:
        logger.warning("Unable to read ODT file: %s", exc)
        return ""
    return extract_openxml_text(xml)


def extract_openxml_text(xml_bytes):
    try:
        root = ElementTree.fromstring(xml_bytes)
    except ElementTree.ParseError:
        return ""

    text_parts = []
    for element in root.iter():
        if element.tag.endswith("}t") and element.text:
            text_parts.append(element.text.strip())

    return " ".join(part for part in text_parts if part).strip()


def split_text(text, chunk_size=1400, overlap=180):
    cleaned = (text or "").strip()
    if not cleaned:
        return []

    chunks = []
    start = 0
    text_length = len(cleaned)

    while start < text_length:
        end = min(text_length, start + chunk_size)
        chunks.append(cleaned[start:end])
        if end >= text_length:
            break
        start = max(0, end - overlap)

    return chunks


def select_relevant_chunks(chunks, query, limit=4):
    if not chunks:
        return []

    query_terms = {term for term in re.findall(r"[a-z0-9]+", (query or "").lower()) if len(term) > 2}
    scored_chunks = []

    for index, chunk in enumerate(chunks):
        lower_chunk = chunk.lower()
        score = sum(1 for term in query_terms if term in lower_chunk)
        scored_chunks.append((score, index, chunk))

    scored_chunks.sort(key=lambda item: (item[0], -item[1]), reverse=True)
    selected = [chunk for score, _, chunk in scored_chunks if score > 0][:limit]

    if selected:
        return selected

    return chunks[:limit]


def build_ai_context(text, max_chars=AI_CONTEXT_LIMIT):
    cleaned = (text or "").strip()
    if len(cleaned) <= max_chars:
        return cleaned

    chunks = split_text(cleaned, chunk_size=1800, overlap=150)
    if len(chunks) <= 10:
        return cleaned[:max_chars]

    middle = len(chunks) // 2
    selected = [
        *chunks[:3],
        *chunks[max(0, middle - 1) : middle + 2],
        *chunks[-3:],
    ]
    return "\n\n[DOCUMENT SAMPLE]\n\n".join(selected)[:max_chars]


def extract_formulas(text, limit=12):
    formulas = []
    patterns = (
        r"\b[A-Za-z][A-Za-z0-9_()^²+-]{0,20}\s*=\s*[^\n;]{2,120}",
        r"\b[A-Z][A-Za-z0-9-]*(?:\s*\([^)]*\))?\s*=\s*[^\n;]{2,120}",
    )
    for pattern in patterns:
        for match in re.findall(pattern, text or ""):
            formula = re.sub(r"\s+", " ", match).strip(" .,:;")
            if formula and formula not in formulas:
                formulas.append(formula)
            if len(formulas) >= limit:
                return formulas
    return formulas


def _split_sentences(text):
    normalized = re.sub(r"\s+", " ", text or "").strip()
    return [sentence.strip() for sentence in re.split(r"(?<=[.!?])\s+", normalized) if sentence.strip()]


def _clean_study_sentence(sentence):
    cleaned = re.sub(r"\s+", " ", sentence or "").strip(" -:;,.")
    cleaned = re.sub(r"https?://\S+|www\.\S+|[\w.+-]+@[\w-]+\.[\w.-]+", "", cleaned)
    cleaned = re.sub(r"\b\+?\d[\d\s-]{7,}\d\b", "", cleaned)
    return re.sub(r"\s+", " ", cleaned).strip(" -:;,.")


def _clean_document_text(text):
    cleaned = re.sub(r"https?://\S+|www\.\S+|[\w.+-]+@[\w-]+\.[\w.-]+", " ", text or "")
    cleaned = re.sub(r"\b\+?\d[\d\s-]{7,}\d\b", " ", cleaned)
    return re.sub(r"\s+", " ", cleaned).strip()


def _candidate_study_sentences(text, limit=12):
    candidates = []

    for sentence in _split_sentences(text):
        cleaned = _clean_study_sentence(sentence)
        lower = cleaned.lower()
        words = re.findall(r"[A-Za-z][A-Za-z+-]*", cleaned)
        if len(words) < 8 or len(cleaned) < 55:
            continue
        if any(term in lower for term in ("email", "phone", "contact", "address", "http", "www")):
            continue

        score = len(set(word.lower() for word in words))
        if any(term in lower for term in ("because", "therefore", "used", "system", "process", "method", "feature", "formula", "value", "metric", "model", "calculate", "compare", "evaluate", "database", "api")):
            score += 10
        candidates.append((score, cleaned))

    candidates.sort(key=lambda item: item[0], reverse=True)
    unique = []
    seen = set()
    for _, sentence in candidates:
        key = sentence[:80].lower()
        if key in seen:
            continue
        seen.add(key)
        unique.append(sentence)
        if len(unique) >= limit:
            break

    if unique:
        return unique

    return [_clean_study_sentence(sentence) for sentence in _split_sentences(text)[:limit] if _clean_study_sentence(sentence)]


def detect_topics(text, top_n=8):
    stopwords = {
        "about",
        "after",
        "also",
        "because",
        "before",
        "between",
        "chapter",
        "course",
        "data",
        "document",
        "example",
        "from",
        "have",
        "lecture",
        "model",
        "more",
        "other",
        "section",
        "should",
        "student",
        "system",
        "than",
        "that",
        "their",
        "there",
        "these",
        "this",
        "using",
        "value",
        "when",
        "which",
        "with",
    }
    counts = {}
    for word in re.findall(r"[A-Za-z][A-Za-z0-9+-]{3,}", text or ""):
        lower = word.lower()
        if lower in stopwords or lower in TOPIC_STOPWORDS:
            continue
        counts[lower] = counts.get(lower, 0) + 1
    ranked = sorted(counts.items(), key=lambda item: (-item[1], item[0]))
    return [word.title() for word, _ in ranked[:top_n]]


def _keywords_from_text(text, limit=4):
    return detect_topics(text, top_n=limit)


def _topic_from_sentence(sentence):
    topics = detect_topics(sentence, top_n=3)
    return " / ".join(topics) if topics else "Core Concept"


def _diverse_study_sentences(text, limit=10):
    candidates = _candidate_study_sentences(text, limit=max(24, limit * 3))
    selected = []
    seen_topics = set()

    for sentence in candidates:
        topic = _topic_from_sentence(sentence)
        if topic in seen_topics:
            continue
        selected.append(sentence)
        seen_topics.add(topic)
        if len(selected) >= limit:
            return selected

    for sentence in candidates:
        if sentence not in selected:
            selected.append(sentence)
        if len(selected) >= limit:
            break

    return selected


def _question_from_sentence(sentence, index, mode="short"):
    lower = sentence.lower()
    topic = _topic_from_sentence(sentence)
    if "formula" in lower or "=" in sentence:
        return f"Which formula or calculation rule is associated with {topic}?"
    if "advantage" in lower or "benefit" in lower:
        return f"Which benefit of {topic} is supported by the material?"
    if "compare" in lower or "different" in lower:
        return f"How should {topic} be compared or interpreted?"
    if mode == "mcq":
        return f"Which statement about {topic} is best supported by the material?"
    return f"What should you remember about {topic}?"


def _compact_quiz_option(text, max_length=180):
    option = _clean_study_sentence(text)
    if len(option) <= max_length:
        return option

    shortened = option[:max_length].rsplit(" ", 1)[0].strip()
    return shortened or option[:max_length].strip()


def _distractors_for_sentence(sentence, context_sentences=None):
    correct_key = _compact_quiz_option(sentence).lower()
    topic_terms = set(re.findall(r"[a-z][a-z0-9+-]{3,}", _topic_from_sentence(sentence).lower()))
    distractors = []

    for candidate in context_sentences or []:
        option = _compact_quiz_option(candidate)
        if not option or option.lower() == correct_key or option in distractors:
            continue

        candidate_terms = set(re.findall(r"[a-z][a-z0-9+-]{3,}", option.lower()))
        if topic_terms and not topic_terms.intersection(candidate_terms):
            continue
        distractors.append(option)
        if len(distractors) == 3:
            return distractors

    for candidate in context_sentences or []:
        option = _compact_quiz_option(candidate)
        if option and option.lower() != correct_key and option not in distractors:
            distractors.append(option)
        if len(distractors) == 3:
            return distractors

    topic = _topic_from_sentence(sentence)
    backup_options = [
        f"{topic} is mainly used for a different step or condition than the one asked here.",
        f"{topic} describes another related idea, but not the specific point required by this question.",
        f"{topic} applies only in a narrower case than the document's supported answer.",
    ]
    return [*distractors, *backup_options][:3]


def _section_content_from_sentences(sentences, limit=5):
    lines = []
    for sentence in sentences[:limit]:
        cleaned = _clean_study_sentence(sentence)
        if cleaned:
            lines.append(f"- {cleaned}")
    return "\n".join(lines)


def _fallback_summary(document):
    text = document.extracted_text or ""
    sentences = _split_sentences(text)
    if not sentences:
        raise ValueError("This file does not contain extractable text. Try a text-based PDF, DOCX, PPTX, XLSX, TXT, Markdown, CSV, JSON, XML, HTML, RTF, or ODT file.")

    summary = []
    topics = detect_topics(text, top_n=8)
    overview = _section_content_from_sentences(sentences, limit=3)
    if overview:
        overview_title = f"Overview of {document.title}" if document.title else "Document Overview"
        summary.append({"section": overview_title, "content": overview[:1000]})

    used_indexes = set()
    for topic in topics:
        topic_terms = {term.lower() for term in re.findall(r"[A-Za-z][A-Za-z0-9+-]{2,}", topic)}
        if not topic_terms:
            continue
        matches = []
        for index, sentence in enumerate(sentences):
            lower = sentence.lower()
            if index in used_indexes:
                continue
            if any(term in lower for term in topic_terms):
                matches.append(sentence)
                used_indexes.add(index)
            if len(matches) >= 5:
                break
        content = _section_content_from_sentences(matches, limit=5)
        if content:
            summary.append({"section": topic, "content": content[:1100]})
        if len(summary) >= 8:
            break

    if len(summary) < 5:
        for chunk in split_text(text, chunk_size=1500, overlap=100)[:6]:
            chunk_sentences = _candidate_study_sentences(chunk, limit=4)
            if not chunk_sentences:
                continue
            section = _topic_from_sentence(" ".join(chunk_sentences[:2]))
            if any(item["section"] == section for item in summary):
                continue
            content = _section_content_from_sentences(chunk_sentences, limit=4)
            if content:
                summary.append({"section": section, "content": content[:1000]})
            if len(summary) >= 8:
                break

    return summary[:8]


def _fallback_flashcards(document):
    text = document.extracted_text or ""
    sentences = _diverse_study_sentences(text, limit=10)
    if not sentences:
        raise ValueError("This file does not contain extractable text. Try a text-based PDF, DOCX, PPTX, XLSX, TXT, Markdown, CSV, JSON, XML, HTML, RTF, or ODT file.")

    cards = []
    for index, sentence in enumerate(sentences[:10]):
        topic = _topic_from_sentence(sentence)
        answer_text = sentence[:320].strip()
        explanation = f"This answer is drawn from the document sentence about {topic}."
        if len(answer_text) > 120:
            explanation = f"The document says: '{answer_text[:120].rstrip()}'... which supports this answer about {topic}."
        cards.append(
            {
                "id": index + 1,
                "question": _question_from_sentence(sentence, index),
                "answer": answer_text,
                "explanation": explanation,
            }
        )
    return cards


def _fallback_multiple_choice_quiz(document):
    text = document.extracted_text or ""
    sentences = _diverse_study_sentences(text, limit=10)
    if not sentences:
        raise ValueError("This file does not contain extractable text. Try a text-based PDF, DOCX, PPTX, XLSX, TXT, Markdown, CSV, JSON, XML, HTML, RTF, or ODT file.")

    quiz_items = []
    for index, sentence in enumerate(sentences[:10]):
        correct = _compact_quiz_option(sentence)
        options = [correct, *_distractors_for_sentence(sentence, sentences)]
        explanation = f"The answer is drawn from the document sentence: '{sentence[:140].rstrip()}'."
        quiz_items.append(
            {
                "id": index + 1,
                "question": _question_from_sentence(sentence, index, mode="mcq"),
                "options": _shuffled_options(options, correct),
                "answer": correct,
                "explanation": explanation,
            }
        )
    return quiz_items


def _fallback_answer(document, question, history=None):
    text = document.extracted_text or ""
    sentences = _split_sentences(text)
    if not sentences:
        raise ValueError("This file does not contain extractable text. Try a text-based PDF, DOCX, PPTX, XLSX, TXT, Markdown, CSV, JSON, XML, HTML, RTF, or ODT file.")

    normalized_question = (question or "").lower()
    overview_terms = ("what does", "what is this", "what this", "about", "shows", "summarize", "summary", "overview")
    is_overview_question = any(term in normalized_question for term in overview_terms)

    if is_overview_question:
        overview_sentences = _candidate_study_sentences(_clean_document_text(text), limit=5)
        if not overview_sentences:
            overview_sentences = [_clean_study_sentence(sentence) for sentence in sentences[:5] if _clean_study_sentence(sentence)]
        return "Based on the uploaded document, here is a concise overview:\n\n" + "\n".join(f"- {sentence}" for sentence in overview_sentences[:5])

    chunks = split_text(text)
    context_chunks = select_relevant_chunks(chunks, question, limit=3)
    matched = ""
    question_terms = {term for term in re.findall(r"[a-z0-9]{4,}", normalized_question)}
    for chunk in context_chunks:
        for sentence in _split_sentences(chunk):
            lower = sentence.lower()
            if len(sentence) > 40 and any(term in lower for term in question_terms):
                matched = sentence.strip()
                break
        if matched:
            break

    if matched:
        return f"From the selected document: {matched}"

    context_sentences = []
    for chunk in context_chunks:
        context_sentences.extend(_split_sentences(chunk))

    concise_context = " ".join(context_sentences[:5]).strip()
    if not concise_context:
        concise_context = " ".join(sentences[:5]).strip()

    if concise_context:
        return (
            "Based on the selected document, the closest matching text is:\n\n"
            f"{concise_context[:900]}"
        )

    return "I could not find a strong answer in the extracted document text. Please try a different question or use a more specific phrase."


def _clean_json_text(text):
    if not text:
        return text

    cleaned = text.strip()
    if cleaned.startswith("```"):
        cleaned = re.sub(r"^```(?:json)?", "", cleaned, flags=re.IGNORECASE).strip()
        cleaned = re.sub(r"```$", "", cleaned).strip()

    start = cleaned.find("[")
    start_object = cleaned.find("{")
    if start == -1 or (start_object != -1 and start_object < start):
        start = start_object

    if start != -1:
        end = cleaned.rfind("]") if cleaned[start] == "[" else cleaned.rfind("}")
        if end != -1 and end > start:
            cleaned = cleaned[start : end + 1]

    return cleaned


class CloudAIService:
    def __init__(self):
        self.groq_key = os.getenv("GROQ_API_KEY", "").strip()
        self.openrouter_key = os.getenv("OPENROUTER_API_KEY", "").strip()
        self.active_provider = str(setting_value("active_ai_provider", "") or "").strip().lower()
        self.active_model = str(setting_value("active_ai_model", "") or "").strip()

    def _post_json(self, url, payload, headers=None, timeout=90):
        data = json.dumps(payload).encode("utf-8")
        request_headers = {
            "Accept": "application/json",
            "Content-Type": "application/json",
            "User-Agent": "StudentAssistant/1.0",
        }
        if headers:
            request_headers.update(headers)

        req = urllib_request.Request(url, data=data, headers=request_headers, method="POST")

        try:
            with urllib_request.urlopen(req, timeout=timeout) as response:
                return json.loads(response.read().decode("utf-8"))
        except urllib_error.HTTPError as exc:
            error_body = exc.read().decode("utf-8", errors="ignore")
            raise RuntimeError(f"AI provider returned HTTP {exc.code}: {error_body[:400]}") from exc
        except urllib_error.URLError as exc:
            raise RuntimeError(f"Unable to reach AI provider: {exc.reason}") from exc

    def _call_groq(self, prompt, system_prompt=None, json_mode=False, temperature=0.2):
        if not self.groq_key:
            raise RuntimeError("Groq API key is not configured.")

        messages = []
        if system_prompt:
            messages.append({"role": "system", "content": system_prompt})
        messages.append({"role": "user", "content": prompt})

        payload = {
            "model": self.active_model or GROQ_MODEL,
            "messages": messages,
            "temperature": temperature,
        }

        endpoint = "https://api.groq.com/openai/v1/chat/completions"
        response = self._post_json(endpoint, payload, headers={"Authorization": f"Bearer {self.groq_key}"})

        choices = response.get("choices") or []
        if not choices:
            raise RuntimeError("Groq returned an empty response.")

        message = choices[0].get("message", {})
        text = message.get("content", "")
        if not text:
            raise RuntimeError("Groq response did not contain text.")

        return text.strip()

    def _call_openrouter_model(self, model, prompt, system_prompt=None, json_mode=False, temperature=0.2):
        if not self.openrouter_key:
            raise RuntimeError("OpenRouter API key is not configured.")

        messages = []
        if system_prompt:
            messages.append({"role": "system", "content": system_prompt})
        messages.append({"role": "user", "content": prompt})

        payload = {
            "model": model,
            "messages": messages,
            "temperature": temperature,
        }

        endpoint = "https://openrouter.ai/api/v1/chat/completions"
        response = self._post_json(
            endpoint,
            payload,
            headers={
                "Authorization": f"Bearer {self.openrouter_key}",
                "HTTP-Referer": os.getenv("OPENROUTER_SITE_URL", "http://localhost:5173"),
                "X-Title": os.getenv("OPENROUTER_APP_NAME", "StudentAssistant"),
            },
        )

        choices = response.get("choices") or []
        if not choices:
            raise RuntimeError(f"OpenRouter model {model} returned an empty response.")

        message = choices[0].get("message", {})
        text = message.get("content", "")
        if not text:
            raise RuntimeError(f"OpenRouter model {model} response did not contain text.")

        return text.strip()

    def _call_openrouter(self, prompt, system_prompt=None, json_mode=False, temperature=0.2):
        errors = []
        models = [self.active_model] if self.active_model else OPENROUTER_MODELS
        for model in models:
            for attempt in range(AI_RETRY_ATTEMPTS):
                try:
                    return self._call_openrouter_model(
                        model,
                        prompt,
                        system_prompt=system_prompt,
                        json_mode=json_mode,
                        temperature=temperature,
                    )
                except Exception as exc:
                    errors.append(f"{model} attempt {attempt + 1}: {exc}")
                    logger.warning("OpenRouter model %s failed on attempt %s: %s", model, attempt + 1, exc)

        if not models:
            raise RuntimeError("OpenRouter models are not configured.")
        raise RuntimeError("All OpenRouter models failed. " + " | ".join(errors[-2:]))

    def generate(self, prompt, system_prompt=None, json_mode=False, temperature=0.2):
        errors = []
        if self.active_provider in {"groq", "grok"}:
            providers = (("Groq", self._call_groq),)
        elif self.active_provider in {"openrouter", "open_router"}:
            providers = (("OpenRouter", self._call_openrouter),)
        else:
            providers = (("Groq", self._call_groq), ("OpenRouter", self._call_openrouter))
        for provider_name, provider in providers:
            for attempt in range(AI_RETRY_ATTEMPTS):
                try:
                    return provider(prompt, system_prompt=system_prompt, json_mode=json_mode, temperature=temperature)
                except Exception as exc:
                    errors.append(f"{provider_name} attempt {attempt + 1}: {exc}")
                    logger.warning("%s failed on attempt %s: %s", provider_name, attempt + 1, exc)
        raise RuntimeError("All AI providers failed. " + " | ".join(errors[-2:]))


def _parse_json_object(raw_text):
    cleaned = _clean_json_text(raw_text)
    try:
        parsed = json.loads(cleaned)
    except Exception:
        parsed = ast.literal_eval(cleaned)
    if isinstance(parsed, dict):
        return parsed
    raise ValueError("AI response did not match the expected JSON object format.")


def classify_document(document):
    text = document.extracted_text or ""
    prompt = (
        "Classify this document into exactly one of these types:\n"
        "Academic Notes, Textbook, Research Paper, Resume, Business Report, Technical Documentation, Manual, Legal Document, Medical Report, Presentation Slides, Other.\n\n"
        'Return ONLY JSON: {"type":"..."}\n\n'
        f"Document title: {document.title}\n\n"
        f"Document text:\n{build_ai_context(text, max_chars=AI_TOPIC_CONTEXT_LIMIT)}"
    )
    try:
        raw = CloudAIService().generate(prompt, system_prompt="Return clean JSON only.", json_mode=True, temperature=0.1)
        classification = str(_parse_json_object(raw).get("type") or "").strip()
        return classification or "Other"
    except Exception as exc:
        logger.warning("AI document classification failed, using fallback: %s", exc)
        lowered = text.lower()
        if any(term in lowered for term in ("abstract", "methodology", "results", "references")):
            return "Research Paper"
        if any(term in lowered for term in ("slide", "agenda", "presentation")):
            return "Presentation Slides"
        if any(term in lowered for term in ("chapter", "exercise", "definition", "formula")):
            return "Academic Notes"
        return "Other"


def extract_document_topics(document, top_n=8):
    text = document.extracted_text or ""
    prompt = (
        "Read this document and identify the most important topics discussed.\n"
        "Return only JSON with 5-10 concise noun-phrase topics.\n"
        "Avoid generic words like Introduction, Chapter, Student, Figure, Learning, Document.\n\n"
        'Return ONLY JSON: {"topics":["Topic 1","Topic 2"]}\n\n'
        f"Document title: {document.title}\n\n"
        f"Document text:\n{build_ai_context(text, max_chars=AI_TOPIC_CONTEXT_LIMIT)}"
    )
    try:
        raw = CloudAIService().generate(prompt, system_prompt="Return clean JSON only.", json_mode=True, temperature=0.1)
        topics = _parse_json_object(raw).get("topics") or []
        topics = [str(topic).strip() for topic in topics if str(topic).strip()]
        if topics:
            return topics[:top_n]
    except Exception as exc:
        logger.warning("AI topic extraction failed, using fallback: %s", exc)
    return detect_topics(text, top_n=top_n)


def build_topic_context(text, topics, max_chars=AI_CONTEXT_LIMIT):
    base_context = build_ai_context(text, max_chars=max_chars)
    chunks = split_text(text, chunk_size=1800, overlap=150)
    if not chunks or not topics:
        return base_context

    sections = []
    for topic in topics[:8]:
        relevant = select_relevant_chunks(chunks, topic, limit=2)
        if relevant:
            sections.append(f"Topic: {topic}\n" + "\n\n".join(relevant))

    topic_context = "\n\n".join(sections).strip()
    if not topic_context:
        return base_context

    return (topic_context + "\n\n[SMART DOCUMENT SAMPLE]\n\n" + base_context)[:max_chars]


def _summary_context_for_generation(document):
    existing_summary = getattr(document, "summary_data", None) or []
    if existing_summary:
        return json.dumps(existing_summary, ensure_ascii=False)
    try:
        return json.dumps(generate_summary(document), ensure_ascii=False)
    except Exception as exc:
        logger.warning("Unable to build summary context, using sampled text: %s", exc)
        return build_ai_context(document.extracted_text or "")


def _build_summary_prompt(document):
    text = document.extracted_text or ""
    if not text.strip():
        raise ValueError("This file does not contain extractable text. Try a text-based PDF, DOCX, PPTX, XLSX, TXT, Markdown, CSV, JSON, XML, HTML, RTF, or ODT file.")

    topics = extract_document_topics(document)
    classification = classify_document(document)
    excerpt = build_topic_context(text, topics)
    formulas = extract_formulas(text)
    return (
        "You are an intelligent study assistant. Create a summary directly from the provided file content and avoid any generic template language.\n"
        "Use only the exact concepts, formulas, examples, comparisons, and statements found in the document.\n"
        "If the source material contains formulas, numbers, or clearly named concepts, preserve those exactly in the summary.\n\n"
        f"Document title: {document.title}\n"
        f"Course or subject label: {document.course or 'General'}\n"
        f"Document type: {classification}\n"
        f"Detected topics: {', '.join(topics) if topics else 'Automatic document topic detection applied'}\n"
        f"Detected formulas or equations: {json.dumps(formulas[:8]) if formulas else 'None detected'}\n\n"
        "Summary requirements:\n"
        "- Return 5 to 8 sections.\n"
        "- Each section should contain 3 to 6 bullets.\n"
        "- Section titles must be specific and derived from the document.\n"
        "- Use concise, student-friendly phrasing anchored in the source text.\n"
        "- Do not invent facts or add unrelated domain knowledge.\n"
        "- Do not use section titles such as Summary, Overview, Key Points, or Main Ideas.\n\n"
        "Return ONLY JSON with the exact schema:\n"
        '[{"section":"...","content":"..."}]\n\n'
        f"Document text:\n{excerpt}"
    )


def _build_flashcard_prompt(document):
    text = document.extracted_text or ""
    if not text.strip():
        raise ValueError("This file does not contain extractable text. Try a text-based PDF, DOCX, PPTX, XLSX, TXT, Markdown, CSV, JSON, XML, HTML, RTF, or ODT file.")

    summary = _summary_context_for_generation(document)
    formulas = extract_formulas(text)
    return (
        "You are an advanced study companion generating high-quality flashcards from a single document.\n"
        "Use the document's actual content, exact wording, and formulas when available.\n"
        "Avoid any generic or broad study questions that are not directly supported by the file.\n\n"
        "Create exactly 10 flashcards.\n"
        "Each flashcard should cover a unique concept, fact, formula, example, or comparison from the document.\n"
        "At least 4 flashcards should include a specific fact, numeric detail, formula, example, or calculation from the file.\n"
        "Ask questions in a way that makes the student recall the precise detail or idea from the file.\n"
        "Provide clear answers with exact statements, definitions, or formulas from the source material.\n"
        "Add a brief explanation showing why the answer is correct based on the document.\n"
        "Do not invent any facts, and do not use generic question templates such as 'What should you remember about ...?'.\n\n"
        "Return ONLY JSON with this schema exactly:\n"
        '[{"id":1,"question":"...","answer":"...","explanation":"..."}]\n\n'
        f"Document title: {document.title}\n"
        f"Detected important formulas or equations: {json.dumps(formulas[:10]) if formulas else 'None detected'}\n\n"
        f"Study summary:\n{summary}\n\n"
        f"Document context:\n{build_ai_context(text)}"
    )


def _build_quiz_prompt(document):
    return _build_flashcard_prompt(document)


def _build_multiple_choice_prompt(document):
    text = document.extracted_text or ""
    if not text.strip():
        raise ValueError("This file does not contain extractable text. Try a text-based PDF, DOCX, PPTX, XLSX, TXT, Markdown, CSV, JSON, XML, HTML, RTF, or ODT file.")

    summary = _summary_context_for_generation(document)
    formulas = extract_formulas(text)
    return (
        "You are an advanced study companion generating high-quality multiple choice questions from a single document.\n"
        "Use the document's actual content, exact wording, and formulas when available.\n"
        "Avoid any generic or broad study questions that are not directly supported by the file.\n\n"
        "Create exactly 10 multiple choice questions.\n"
        "Each question should cover a unique concept, fact, formula, example, comparison, or procedure from the document.\n"
        "At least 4 questions should include a specific fact, numeric detail, formula, example, or calculation from the file.\n"
        "Use clear and specific wording that the student could answer from the source text.\n"
        "Provide concise answer explanations that directly reference the document.\n"
        "Do not invent any facts, and do not use generic question templates or unrelated phrasing.\n\n"
        "Return ONLY JSON with this schema exactly:\n"
        '[{"id":1,"question":"...","options":["","","",""],"answer":"...","explanation":"..."}]\n\n'
        f"Document title: {document.title}\n"
        f"Detected important formulas or equations: {json.dumps(formulas[:10]) if formulas else 'None detected'}\n\n"
        f"Study summary:\n{summary}\n\n"
        f"Document context:\n{build_ai_context(text)}"
    )

def _build_chat_prompt(document, question, context_chunks, history=None):
    history_lines = []
    for item in (history or [])[-8:]:
        role = item.get("role", "user")
        message = item.get("message", "")
        if message:
            history_lines.append(f"{role.upper()}: {message}")

    context = "\n\n".join(context_chunks) if context_chunks else "No direct context found in the uploaded document."
    history_text = "\n".join(history_lines) if history_lines else "No prior conversation."

    return (
        "You are an advanced AI study companion helping a student understand their study materials.\n"
        "The student has uploaded a document for study. Answer their questions intelligently and conversationally.\n\n"
        "Guidelines:\n"
        "- First, try to answer from the provided document context when it directly applies.\n"
        "- For requests like 'give an example', 'explain this concept', 'elaborate on', or 'what does this mean', you can provide examples and explanations beyond the document.\n"
        "- If the student asks a general question about the topic, you can answer it using broader knowledge, not just the document.\n"
        "- When providing additional context, always relate it back to what the student is studying.\n"
        "- Be conversational, helpful, and encouraging.\n"
        "- If something is directly in the document, mention it: 'The document says...' or 'In your material...'.\n"
        "- Do not make up specific facts from the document that you don't see in the context.\n\n"
        "Formatting requirements:\n"
        "- Use clean Markdown with short sections, bullets, or numbered steps when useful.\n"
        "- Put a blank line before each heading, list, or table.\n"
        "- Use headings like ## Topic or ### Next Steps only when they help readability.\n"
        "- Use valid Markdown tables only when comparing multiple items, with a header row and separator row.\n"
        "- Do not cram headings, bullets, and tables into one paragraph.\n\n"
        f"Student's document: {document.title}\n"
        f"Subject/Course: {document.course or 'General'}\n\n"
        f"Conversation history:\n{history_text}\n\n"
        f"Relevant excerpt from the document (if available):\n{context}\n\n"
        f"Student's question: {question}"
    )


def generate_valid_json_array(prompt, possible_keys, validator, temperature=0.2):
    last_error = None
    service = CloudAIService()
    retry_prompt = prompt
    for attempt in range(3):
        try:
            raw = service.generate(
                retry_prompt,
                system_prompt="Return clean JSON only. No markdown. No commentary.",
                json_mode=True,
                temperature=temperature,
            )
            parsed = parse_json_array(raw, possible_keys=possible_keys)
            if not isinstance(parsed, list):
                raise ValueError("AI response must be a JSON array.")
            validator(parsed)
            return parsed
        except Exception as exc:
            last_error = exc
            logger.warning("AI JSON generation validation failed on attempt %s: %s", attempt + 1, exc)
            retry_prompt = (
                prompt
                + "\n\nYour previous response failed validation. Return ONLY valid JSON matching the exact schema, with no markdown or extra text."
            )
    raise ValueError(f"AI response failed validation after retries: {last_error}")


def generate_summary(document):
    result = generate_summary_with_source(document)
    return result["items"]


def generate_summary_with_source(document):
    prompt = _build_summary_prompt(document)
    try:
        parsed = generate_valid_json_array(
            prompt,
            possible_keys=["summary", "sections", "items"],
            validator=lambda items: normalize_summary_sections(items, document),
            temperature=0.2,
        )
        return {"items": normalize_summary_sections(parsed, document), "source": "ai"}
    except Exception as exc:
        logger.warning("Cloud summary failed, using local fallback: %s", exc)
        return {"items": _fallback_summary(document), "source": "fallback"}


def generate_flashcards(document):
    result = generate_flashcards_with_source(document)
    return result["items"]


def generate_flashcards_with_source(document):
    prompt = _build_flashcard_prompt(document)
    try:
        parsed = generate_valid_json_array(
            prompt,
            possible_keys=["flashcards", "cards", "questions", "items"],
            validator=normalize_flashcards,
            temperature=0.25,
        )
        return {"items": normalize_flashcards(parsed), "source": "ai"}
    except Exception as exc:
        logger.warning("Cloud flashcard generation failed, using local fallback: %s", exc)
        return {"items": _fallback_flashcards(document), "source": "fallback"}


def generate_multiple_choice_quiz(document):
    result = generate_multiple_choice_quiz_with_source(document)
    return result["items"]


def generate_multiple_choice_quiz_with_source(document):
    prompt = _build_multiple_choice_prompt(document)
    try:
        parsed = generate_valid_json_array(
            prompt,
            possible_keys=["quiz", "questions", "items"],
            validator=normalize_quiz_items,
            temperature=0.25,
        )
        return {"items": normalize_quiz_items(parsed), "source": "ai"}
    except Exception as exc:
        logger.warning("Cloud MCQ quiz generation failed, using local fallback: %s", exc)
        return {"items": _fallback_multiple_choice_quiz(document), "source": "fallback"}


def answer_question(document, question, history=None):
    if not (document.extracted_text or "").strip():
        raise ValueError("This file does not contain extractable text. Try a text-based PDF, DOCX, PPTX, XLSX, TXT, Markdown, CSV, JSON, XML, HTML, RTF, or ODT file.")

    chunks = split_text(document.extracted_text)
    context_chunks = select_relevant_chunks(chunks, question, limit=4)
    prompt = _build_chat_prompt(document, question, context_chunks, history=history)
    try:
        response = CloudAIService().generate(
            prompt,
            system_prompt="You are a precise study assistant.",
            json_mode=False,
            temperature=0.2,
        )
        return {
            "message": response,
            "is_ai_generated": True,
            "source": "ai",
        }
    except Exception as exc:
        error_msg = str(exc)
        logger.error("Cloud chat API failed: %s", error_msg)
        logger.warning("Using document-based fallback for chat question: %s", question)
        fallback_response = _fallback_answer(document, question, history=history)
        return {
            "message": fallback_response,
            "is_ai_generated": False,
            "source": "api_error" if ("API key" in error_msg or "not configured" in error_msg.lower()) else "fallback",
        }


def parse_json_array(raw_text, possible_keys):
    cleaned = _clean_json_text(raw_text)
    try:
        parsed = json.loads(cleaned)
    except Exception:
        parsed = ast.literal_eval(cleaned)

    if isinstance(parsed, list):
        return parsed

    if isinstance(parsed, dict):
        for key in possible_keys:
            value = parsed.get(key)
            if isinstance(value, list):
                return value

    raise ValueError("AI response did not match the expected JSON format.")


def normalize_summary_sections(items, document=None):
    normalized = []
    for index, item in enumerate(items):
        if not isinstance(item, dict):
            continue
        section = str(item.get("section") or item.get("title") or "").strip()
        content = item.get("content") or item.get("summary") or item.get("notes") or ""
        if isinstance(content, list):
            content = "\n".join(f"- {str(line).strip()}" for line in content if str(line).strip())
        content = str(content).strip()
        if not content:
            continue
        if not section:
            section = _topic_from_sentence(content) or f"Study Point {index + 1}"
        normalized.append(
            {
                "section": section,
                "content": content,
            }
        )

    if not normalized:
        raise ValueError("Summary response did not contain valid sections.")

    if len(normalized) < 3 and document:
        return _fallback_summary(document)

    return normalized[:8]


def normalize_quiz_items(items):
    normalized = []
    seen_questions = set()
    for index, item in enumerate(items):
        if not isinstance(item, dict):
            continue

        options = item.get("options") or []
        if not isinstance(options, list):
            options = []
        options = [str(option).strip() for option in options if str(option).strip()]
        unique_options = []
        for option in options:
            if option not in unique_options:
                unique_options.append(option)
        options = unique_options

        answer = str(item.get("answer") or "").strip()
        if answer and answer not in options:
            options = [answer, *options][:4]
        elif answer and len(options) > 4:
            options = [answer, *[option for option in options if option != answer]][:4]

        if len(options) < 4:
            options = [*options, *_distractors_for_sentence(answer or question_from_item(item))][:4]
        unique_options = []
        for option in options:
            if option and option not in unique_options:
                unique_options.append(option)
        options = unique_options

        if not answer or len(options) != 4 or answer not in options:
            continue

        question = str(item.get("question") or f"Question {index + 1}").strip()
        question_key = re.sub(r"\W+", " ", question.lower()).strip()
        if len(question) <= 15 or question_key in seen_questions:
            continue
        seen_questions.add(question_key)

        normalized.append(
            {
                "id": item.get("id") or index + 1,
                "question": question,
                "options": _shuffled_options(options, answer),
                "answer": answer,
                "explanation": str(item.get("explanation") or "").strip(),
            }
        )

    if len(normalized) < 6:
        raise ValueError("Quiz response did not contain valid multiple choice questions.")

    return normalized


def question_from_item(item):
    return str(item.get("question") or item.get("answer") or "")


def _shuffled_options(options, answer):
    shuffled = list(options)
    random.shuffle(shuffled)
    if len(shuffled) > 1 and shuffled[0] == answer:
        rotate_by = random.randint(1, len(shuffled) - 1)
        shuffled = shuffled[rotate_by:] + shuffled[:rotate_by]
    return shuffled




def normalize_flashcards(items):
    normalized = []
    seen_questions = set()
    for index, item in enumerate(items):
        if not isinstance(item, dict):
            continue
        question = str(item.get("question") or "").strip()
        answer = str(item.get("answer") or "").strip()
        question_key = re.sub(r"\W+", " ", question.lower()).strip()
        if not question or not answer or question_key in seen_questions:
            continue
        if len(question) <= 15:
            continue
        seen_questions.add(question_key)
        normalized.append(
            {
                "id": item.get("id") or index + 1,
                "question": question,
                "answer": answer,
                "explanation": str(item.get("explanation") or "").strip(),
            }
        )

    if len(normalized) < 6:
        raise ValueError("Flashcard response did not contain valid study cards.")

    return normalized[:10]

